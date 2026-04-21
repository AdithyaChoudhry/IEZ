"""
_build_ppt.py
Corporate presentation — Instrumentation EZ + IODB Validation
Manager / MD audience.  Run:  python _build_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (classic corporate blue / steel) ───────────────────────────
NAVY     = RGBColor(0x0D, 0x2B, 0x4E)   # deep navy
BLUE     = RGBColor(0x1B, 0x5E, 0xA7)   # corporate blue
LBLUE    = RGBColor(0xD6, 0xE4, 0xF7)   # light blue fill
MID      = RGBColor(0x26, 0x7A, 0xB5)   # mid blue accents
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W    = RGBColor(0xF4, 0xF7, 0xFB)   # slide background
DGRAY    = RGBColor(0x33, 0x3A, 0x45)   # body text dark
MGRAY    = RGBColor(0x60, 0x6A, 0x78)   # secondary text
LGRAY    = RGBColor(0xE8, 0xEC, 0xF1)   # light rule / card bg
GREEN    = RGBColor(0x1B, 0x7A, 0x4E)
ORANGE   = RGBColor(0xC9, 0x5F, 0x10)
RED      = RGBColor(0xB3, 0x20, 0x20)
GOLD     = RGBColor(0xC8, 0x9A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ═════════════════════════ LOW-LEVEL HELPERS ══════════════════════════════════

def rgb(r, g, b): return RGBColor(r, g, b)

def rect(sl, l, t, w, h, fill, line=False):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if not line: s.line.fill.background()
    else:
        s.line.fill.solid(); s.line.fore_color.rgb = fill
    return s

def txt(sl, text, l, t, w, h, size=14, bold=False, italic=False,
        color=DGRAY, align=PP_ALIGN.LEFT, wrap=True, name="Calibri Light"):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = name
    return tb

def para(tf, text, size=13, bold=False, color=DGRAY,
         align=PP_ALIGN.LEFT, space=5, name="Calibri"):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = name; return p

def slide_base(sl):
    """Off-white background + thin navy left accent bar."""
    rect(sl, 0, 0, 13.33, 7.5, OFF_W)
    rect(sl, 0, 0, 0.07,  7.5, NAVY)

def hdr(sl, title, subtitle=""):
    """Standard header: navy top bar + white title + blue subtitle."""
    rect(sl, 0, 0, 13.33, 1.15, NAVY)
    rect(sl, 0, 1.15, 13.33, 0.04, BLUE)
    txt(sl, title,    0.35, 0.08, 12.5, 0.72, size=26, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT, name="Calibri Light")
    if subtitle:
        txt(sl, subtitle, 0.35, 0.78, 12.5, 0.35, size=13,
            color=rgb(0xA8,0xC8,0xEE), align=PP_ALIGN.LEFT, name="Calibri")

def foot(sl, note="Instrumentation EZ  |  Internal Presentation  |  2026  |  Confidential"):
    rect(sl, 0, 7.28, 13.33, 0.22, NAVY)
    txt(sl, note, 0.35, 7.29, 12.5, 0.2,
        size=8, color=rgb(0x90,0xA8,0xC4), align=PP_ALIGN.CENTER, name="Calibri")

def bullet_box(sl, l, t, w, h, title, items, title_bg=BLUE,
               title_size=13, item_size=12, item_color=DGRAY, prefix="- "):
    """Card with colored header band and bullet list below."""
    rect(sl, l, t, w, 0.42, title_bg)
    txt(sl, title, l+0.12, t+0.05, w-0.24, 0.35,
        size=title_size, bold=True, color=WHITE, name="Calibri")
    body = sl.shapes.add_textbox(
        Inches(l+0.12), Inches(t+0.5), Inches(w-0.24), Inches(h-0.55))
    tf = body.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(5)
        r = p.add_run(); r.text = prefix + item
        r.font.size = Pt(item_size); r.font.color.rgb = item_color
        r.font.name = "Calibri"

def stat_card(sl, l, t, w, h, big, small, desc, color=BLUE):
    rect(sl, l, t, w, h, WHITE)
    rect(sl, l, t, w, 0.07, color)
    txt(sl, big,   l+0.1, t+0.13, w-0.2, h*0.45,
        size=38, bold=True, color=color, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, small, l+0.1, t+h*0.55, w-0.2, 0.38,
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, desc,  l+0.1, t+h*0.72, w-0.2, 0.42,
        size=10, color=MGRAY, align=PP_ALIGN.CENTER, name="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 0, 13.33, 0.08, BLUE)
rect(sl, 0, 7.42, 13.33, 0.08, BLUE)
acc = sl.shapes.add_shape(1, Inches(8.8), Inches(0), Inches(4.53), Inches(7.5))
acc.fill.solid(); acc.fill.fore_color.rgb = rgb(0x12,0x38,0x62)
acc.line.fill.background()

txt(sl, "Instrumentation EZ", 0.65, 1.1, 9.5, 1.05,
    size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT, name="Calibri Light")
txt(sl, "Automated Engineering Document Generation & IODB Validation Platform",
    0.65, 2.25, 8.5, 0.9, size=22, color=rgb(0xA8,0xC8,0xEE),
    align=PP_ALIGN.LEFT, name="Calibri Light")
rect(sl, 0.65, 3.3, 5.8, 0.05, BLUE)
txt(sl, "Presented to: Management & Executive Leadership",
    0.65, 3.5, 8.5, 0.4, size=14, color=rgb(0x88,0xA8,0xCC),
    align=PP_ALIGN.LEFT, name="Calibri")
txt(sl, "April 2026  |  Instrumentation Engineering Team",
    0.65, 3.95, 8.5, 0.38, size=13, color=rgb(0x70,0x90,0xB0),
    align=PP_ALIGN.LEFT, name="Calibri")

modules = ["Instrument List","I/O List","Data Sheet","Cable Schedule",
           "Loop Wiring","IODB Validation"]
for i, m in enumerate(modules):
    x = 0.65 + (i % 3)*2.72
    y = 5.2  + (i // 3)*0.65
    rect(sl, x, y, 2.55, 0.5, rgb(0x1E,0x4A,0x82))
    txt(sl, m, x+0.08, y+0.06, 2.4, 0.38,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")

txt(sl, "CONFIDENTIAL", 9.5, 6.9, 3.5, 0.4,
    size=10, color=rgb(0x60,0x80,0xA0), align=PP_ALIGN.RIGHT, name="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl); hdr(sl, "Agenda", "What we will cover today")

sections = [
    ("01", "Executive Summary",         "Business case and purpose of the platform"),
    ("02", "Problem Statement",         "Current challenges in instrumentation engineering"),
    ("03", "Solution Overview",         "Introduction to Instrumentation EZ"),
    ("04", "System Architecture",       "How the platform is structured and works"),
    ("05", "Module Deep-Dives",         "Instrument List · I/O List · Data Sheet · Cable Schedule · Loop Wiring"),
    ("06", "IODB Validation Engine",    "How validation works, error detection, Excel report"),
    ("07", "Key Benefits & ROI",        "Quantified time savings and quality gains"),
    ("08", "Demo Walkthrough",          "Live demonstration flow"),
    ("09", "Future Roadmap",            "Planned enhancements and scalability"),
    ("10", "Conclusion",                "Business impact and next steps"),
]
for i, (num, title, desc) in enumerate(sections):
    row = i // 2; col = i % 2
    x = 0.35 + col*6.5; y = 1.35 + row*1.15
    rect(sl, x, y, 6.2, 0.98, WHITE)
    rect(sl, x, y, 0.55, 0.98, BLUE)
    txt(sl, num, x+0.04, y+0.18, 0.5, 0.5,
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, title, x+0.65, y+0.05, 5.4, 0.42,
        size=14, bold=True, color=NAVY, name="Calibri")
    txt(sl, desc,  x+0.65, y+0.52, 5.4, 0.38,
        size=11, color=MGRAY, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Executive Summary",
    "Instrumentation EZ — purpose, scope, and business value at a glance")

rect(sl, 0.35, 1.3, 7.8, 5.5, WHITE)
rect(sl, 0.35, 1.3, 7.8, 0.06, BLUE)
body = sl.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(7.5), Inches(5.1))
tf = body.text_frame; tf.word_wrap = True
first = True

paras = [
    ("What is it?",
     "Instrumentation EZ is an internal web-based platform that automates the "
     "generation and validation of engineering documents directly from the "
     "Instrument & Operations Database (IODB)."),
    ("Why does it exist?",
     "Engineering teams spend 60-70% of document-preparation time on manual data "
     "extraction, formatting, and cross-checking — work that is error-prone and "
     "non-value-adding. This platform eliminates that burden."),
    ("What does it do?",
     "It provides 6 integrated modules: Instrument List, I/O List, Data Sheet, "
     "Cable Schedule, Loop Wiring — and an IODB Validation Engine that "
     "automatically detects data quality issues and logs them in a structured Excel report."),
    ("Who benefits?",
     "Instrumentation engineers, project leads, DCS/PLC engineers, procurement, "
     "and QA teams — anyone who consumes or reviews these deliverables."),
    ("Core value:",
     "One upload — five documents generated — zero transcription errors — "
     "full audit trail of IODB data quality."),
]
for heading, body_text in paras:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False; p.space_before = Pt(10)
    r  = p.add_run(); r.text = heading + "  "
    r.font.bold = True; r.font.size = Pt(13)
    r.font.color.rgb = BLUE; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = body_text
    r2.font.size = Pt(12); r2.font.color.rgb = DGRAY; r2.font.name = "Calibri"

kpis = [
    (">95%", "Time Saved",    "vs. manual process", BLUE),
    ("6",    "Modules",       "in one platform",    MID),
    ("~0",   "Manual Errors", "with automation",    GREEN),
    ("100%", "Template Safe", "formats preserved",  rgb(0x6B,0x35,0xAF)),
]
for i, (val, lbl, sub, c) in enumerate(kpis):
    stat_card(sl, 8.42, 1.3 + i*1.38, 4.5, 1.22, val, lbl, sub, c)

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Problem Statement",
    "The hidden cost of manual instrumentation document preparation")

problems = [
    (ORANGE, "Time-Intensive Manual Work",
     ["Copy-pasting data from IODB into every document type",
      "A single Datasheet takes 15-30 minutes to prepare manually",
      "100 instruments = 3-4 weeks of engineering time",
      "Every IODB revision requires full rework of all documents"]),
    (RED,    "High Error Risk",
     ["Transcription errors undetected until client review",
      "Inconsistencies between Instrument List and Cable Schedule",
      "Missing values left blank instead of flagged",
      "No systematic data quality check across the IODB"]),
    (GOLD,   "No Standardisation",
     ["Different engineers produce differently formatted documents",
      "Template versions mismatched across project phases",
      "No traceability between IODB version and output document",
      "Rework caused by late detection of upstream data errors"]),
]
for i, (color, title, bullets) in enumerate(problems):
    x = 0.35 + i*4.3
    rect(sl, x, 1.3, 4.1, 5.5, WHITE)
    rect(sl, x, 1.3, 4.1, 0.07, color)
    txt(sl, title, x+0.15, 1.42, 3.85, 0.52,
        size=13, bold=True, color=NAVY, name="Calibri")
    body = sl.shapes.add_textbox(
        Inches(x+0.15), Inches(2.0), Inches(3.85), Inches(4.6))
    tf = body.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(8)
        r = p.add_run(); r.text = "- " + b
        r.font.size = Pt(13); r.font.color.rgb = MGRAY; r.font.name = "Calibri"

rect(sl, 0.35, 6.95, 12.65, 0.38, NAVY)
txt(sl,
    "Industry data: Manual document preparation accounts for up to 40% of engineering "
    "labour on mid-size projects — the highest non-value-adding activity on site.",
    0.5, 6.97, 12.3, 0.35, size=11,
    color=rgb(0xC5,0xD8,0xF0), align=PP_ALIGN.CENTER, name="Calibri")
foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Solution Overview — Instrumentation EZ",
    "A single platform that turns raw IODB data into finished, formatted engineering documents")

txt(sl,
    "Upload once.  Generate everything.  Validate instantly.",
    0.35, 1.3, 12.65, 0.6, size=22, bold=True,
    color=NAVY, align=PP_ALIGN.CENTER, name="Calibri Light")
rect(sl, 3.0, 2.0, 7.35, 0.04, BLUE)

flow = [
    ("Upload\nIODB", BLUE),
    ("Select\nModule", MID),
    ("Configure\nOptions", rgb(0x1B,0x7A,0x4E)),
    ("Generate", ORANGE),
    ("Download\nOutput", NAVY),
]
box_w, box_h = 1.85, 1.5
start_x = 0.55
for i, (label, color) in enumerate(flow):
    x = start_x + i*2.55
    rect(sl, x, 2.2, box_w, box_h, color)
    txt(sl, label, x+0.05, 2.25, box_w-0.1, box_h-0.1,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    if i < 4:
        txt(sl, "->", x+box_w+0.08, 2.72, 0.52, 0.5,
            size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

M = ["Instrument List","I/O List","Data Sheet","Cable Schedule","Loop Wiring","IODB Validation"]
colors_m = [BLUE, MID, rgb(0x1B,0x7A,0x4E), ORANGE, NAVY, RED]
for i, (m, c) in enumerate(zip(M, colors_m)):
    x = 0.35 + (i%3)*4.22; y = 4.0 + (i//3)*0.72
    rect(sl, x, y, 4.0, 0.58, c)
    txt(sl, m, x+0.12, y+0.1, 3.8, 0.4,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT, name="Calibri")

rect(sl, 0.35, 6.95, 12.65, 0.38, LGRAY)
txt(sl,
    "Built with Python 3.14 · Streamlit · openpyxl · pandas  |  "
    "Runs in any browser — no installation required for end users",
    0.5, 6.97, 12.3, 0.35, size=11, color=MGRAY,
    align=PP_ALIGN.CENTER, name="Calibri")
foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "System Architecture",
    "Clean four-layer design — easy to extend, maintain, and deploy")

layers = [
    (BLUE,  "PRESENTATION LAYER",
     "main.py  —  Streamlit Web Application",
     "Module selector · File upload widgets · Progress indicators · "
     "Preview tables · Download buttons · Session-state file caching"),
    (MID,   "ORCHESTRATION LAYER",
     "utils/template_processor.py  —  Pipeline Coordinator",
     "Reads and parses uploaded files · Invokes the correct generator · "
     "Validates data shapes · Bundles multi-file ZIP outputs"),
    (GREEN, "DATA LAYER",
     "utils/file_handler.py  —  Robust File Ingestion",
     "Auto-detects sheet name & header row · Fuzzy column matching · "
     "Content_Types.xml auto-repair · .xls format detection & rejection"),
    (ORANGE,"GENERATION LAYER",
     "utils/generators.py  —  Document & Validation Engines",
     "5 document generators + 1 IODB Validation engine · "
     "Formula preservation · Drawing/shape injection · TBA safety net"),
]
for i, (c, layer, comp, desc) in enumerate(layers):
    y = 1.28 + i*1.45
    rect(sl, 0.35, y, 12.65, 1.25, WHITE)
    rect(sl, 0.35, y, 1.55,  1.25, c)
    txt(sl, layer, 0.37, y+0.38, 1.52, 0.48,
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, comp,  2.05, y+0.08, 5.5, 0.45,
        size=14, bold=True, color=NAVY, name="Calibri")
    txt(sl, desc,  2.05, y+0.58, 10.7, 0.55,
        size=11, color=MGRAY, name="Calibri")
    if i < 3:
        txt(sl, "v", 6.58, y+1.27, 0.5, 0.32,
            size=14, color=BLUE, align=PP_ALIGN.CENTER)

txt(sl, "INPUT  ->  IODB.xlsx  +  Template.xlsx",
    0.5, 7.02, 6.0, 0.28, size=11, bold=True, color=BLUE)
txt(sl, "OUTPUT  ->  .xlsx / .zip  — Ready for issue",
    6.9, 7.02, 6.0, 0.28, size=11, bold=True,
    color=GREEN, align=PP_ALIGN.RIGHT)
foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODULES OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Platform Modules — At a Glance",
    "Five automated document generators + one validation engine")

mods = [
    (BLUE,  "01  Instrument List",
     "Exports any selected IODB columns as a\nclean, formatted Excel register.",
     "< 5 sec", "3-4 hrs"),
    (MID,   "02  I/O List",
     "Auto-suggests I/O columns; produces\na DCS/PLC-ready signal register.",
     "< 5 sec", "4-6 hrs"),
    (rgb(0x1B,0x7A,0x4E), "03  Data Sheet",
     "Per-tag datasheets from template;\nbatch ZIP export for all tags.",
     "~2 sec/tag","15-30 min/tag"),
    (ORANGE,"04  Cable Schedule",
     "JB-grouped, naturally sorted,\nSPARE-padded cable schedule.",
     "< 30 sec", "1-2 days"),
    (NAVY,  "05  Loop Wiring",
     "Per-tag loop sheets; shapes, lines\nand formulas fully preserved.",
     "~3 sec/tag","30-60 min/tag"),
    (RED,   "06  IODB Validation",
     "Scans entire IODB; logs every error\nwith location + type + description.",
     "< 10 sec", "Manual audit"),
]
for i, (c, title, desc, after, before) in enumerate(mods):
    row = i // 3; col = i % 3
    x = 0.35 + col*4.35; y = 1.28 + row*2.85
    rect(sl, x, y, 4.15, 2.65, WHITE)
    rect(sl, x, y, 4.15, 0.44, c)
    txt(sl, title, x+0.12, y+0.05, 3.92, 0.36,
        size=13, bold=True, color=WHITE, name="Calibri")
    txt(sl, desc,  x+0.12, y+0.52, 3.92, 0.85,
        size=11.5, color=DGRAY, name="Calibri")
    rect(sl, x+0.12, y+1.45, 3.92, 0.03, LGRAY)
    txt(sl, "Before: " + before,
        x+0.12, y+1.55, 1.88, 0.38, size=10, color=RED, bold=True, name="Calibri")
    txt(sl, "After: " + after,
        x+2.1,  y+1.55, 2.0,  0.38, size=10, color=GREEN, bold=True,
        align=PP_ALIGN.RIGHT, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DATA SHEET MODULE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Module Deep-Dive: Data Sheet Generator",
    "Batch-generate fully formatted instrument datasheets — one per tag — from IODB + template")

bullet_box(sl, 0.35, 1.28, 5.9, 2.7, "How It Works", [
    "Upload IODB source + Datasheet template",
    "Engine maps template labels to IODB columns using fuzzy matching",
    "For each selected tag: copy template, fill matched cells",
    "Formula cells are never overwritten — calculated fields stay live",
    "Empty IODB values written as 'TBA' (never left blank)",
    "All datasheets packaged into Datasheets.zip for one-click download",
], title_bg=BLUE)

bullet_box(sl, 6.45, 1.28, 6.5, 2.7, "Smart Matching Engine", [
    "'Line Size (NB)' in template auto-maps to 'line size' in IODB",
    "Works across any template format — no manual column mapping",
    "Header row auto-detected even if data starts at row 3-6",
    "Partial / substring match as fallback for non-standard names",
    "Case-insensitive across all comparisons",
], title_bg=MID)

bullet_box(sl, 0.35, 4.2, 5.9, 2.65, "Business Impact", [
    "100 datasheets in < 4 minutes vs. ~50 hours manually",
    "Zero transcription errors — IODB is the single source of truth",
    "Engineers focus on engineering, not copy-paste",
    "Re-run instantly on any IODB revision",
    "Consistent formatting across every tag and every project",
], title_bg=GREEN)

bullet_box(sl, 6.45, 4.2, 6.5, 2.65, "Supported Use Cases", [
    "IFC (Issued for Construction) datasheet packages",
    "As-built datasheet generation post-commissioning",
    "Bid / tender datasheet packages for vendor inquiry",
    "Internal review packs for engineering sign-off",
    "Client deliverable compilation for project handover",
], title_bg=ORANGE)

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CABLE SCHEDULE MODULE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Module Deep-Dive: Cable Schedule Generator",
    "Junction-box-grouped, naturally sorted, SPARE-padded cable schedule — in under 30 seconds")

bullet_box(sl, 0.35, 1.28, 8.1, 5.55, "Processing Pipeline", [
    "STEP 1 — Load IODB, auto-detect 'Junction Box' and 'Tag No' columns",
    "STEP 2 — Filter out tags with JB = '-' (unassigned / field instruments)",
    "STEP 3 — Group all tags by Junction Box reference number",
    "STEP 4 — Natural-sort each JB group (1, 2, 3 ... 10, not 1, 10, 2)",
    "STEP 5 — Pad each group to 12 rows with SPARE entries (field-change buffer)",
    "STEP 6 — Match IODB column names to template headers dynamically",
    "STEP 7 — Fill template rows; preserve all borders, merged cells, row heights",
    "STEP 8 — Adjust formula row references when rows are inserted",
    "STEP 9 — Compact output by removing blank spacer rows",
    "STEP 10 — Apply Arial 12 centre-wrap formatting to all data cells",
], title_bg=BLUE, item_size=12)

bullet_box(sl, 8.65, 1.28, 4.3, 2.65, "Engineering Value", [
    "Mirrors physical wiring layout in the field",
    "SPARE rows allow field changes without rework",
    "Formula breakage on row insert — solved",
    "Multi-JB project: all scheduled in one run",
    "Output: IFC-ready Cable_Schedule.xlsx",
], title_bg=ORANGE)

bullet_box(sl, 8.65, 4.15, 4.3, 2.68, "Time Comparison", [
    "Manual:  1-2 days per revision",
    "With tool:  < 30 seconds",
    "Re-run cost after IODB change:",
    "  Manual:  Full rework (hours)",
    "  With tool:  Re-click (seconds)",
], title_bg=GREEN)

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LOOP WIRING MODULE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Module Deep-Dive: Loop Wiring Generator",
    "Per-tag loop sheets with shapes, connectors, and formulas fully preserved")

bullet_box(sl, 0.35, 1.28, 6.1, 5.55, "How It Works", [
    "Upload Loop Wiring Input file (tag list + instrument data)",
    "Locate 'AI Instr' sheet in template (tolerant name matching)",
    "Read column headers from Row 18, columns W to AR (22 fields)",
    "For each tag: copy_worksheet() — full template duplication",
    "Fill Row 19 (W to AR) with matched values from input file",
    "Formula cells never overwritten — all live calculations preserved",
    "Empty values written as 'TBA' — no blank cells in output",
    "Sheet named 'Loop_<TagNumber>' for clear navigation",
    "Tolerant matching: 'AI-INST ', 'ai instr', 'AI Instr' all work",
    "Output: Loop_Wiring.xlsx with one tab per instrument",
], title_bg=BLUE)

bullet_box(sl, 6.7, 1.28, 6.25, 2.7, "Drawing & Shape Preservation", [
    "openpyxl copy_worksheet() preserves native shapes & styles",
    "For complex drawings (connectors, P&ID symbols, text boxes)",
    "that openpyxl cannot round-trip: ZIP-level XML injection",
    "Drawing XML copied directly into output ZIP after save",
    "Every generated sheet retains the original template drawings",
], title_bg=MID)

bullet_box(sl, 6.7, 4.2, 6.25, 2.65, "Business Impact", [
    "100-tag loop wiring set: ~5 mins vs. ~80 hrs manually",
    "No risk of missing fields — TBA ensures full coverage",
    "Template revision: re-run with new template — done",
    "Consistent sheet naming aids document control",
    "Directly usable as field loop folder / commissioning pack",
], title_bg=GREEN)

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IODB VALIDATION: INTRO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "IODB Validation Engine",
    "Automated data quality analysis — catch every error before it reaches a drawing or document")

rect(sl, 0.35, 1.28, 12.65, 0.75, NAVY)
txt(sl,
    "The IODB is the single source of truth for all engineering documents. "
    "A single error propagates into every output — Validation stops that.",
    0.5, 1.33, 12.3, 0.6, size=14, color=WHITE,
    align=PP_ALIGN.CENTER, name="Calibri")

pillars = [
    (BLUE,   "Analyse",
     "Reads the entire IODB and inspects every row and column "
     "for completeness, format compliance, and cross-field consistency."),
    (ORANGE, "Detect",
     "Applies a rule library to identify: missing mandatory fields, "
     "duplicate tag numbers, invalid signal types, format violations, "
     "and cross-reference mismatches."),
    (GREEN,  "Report",
     "Generates a structured Excel log with exact cell location "
     "(sheet, row, column), error type, severity level, and "
     "a plain-English description of what is wrong and why."),
]
for i, (c, title, desc) in enumerate(pillars):
    x = 0.35 + i*4.3
    rect(sl, x, 2.2, 4.1, 4.45, WHITE)
    rect(sl, x, 2.2, 4.1, 0.07, c)
    txt(sl, title, x+0.15, 3.6, 3.85, 0.42,
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, desc, x+0.15, 4.15, 3.85, 2.3,
        size=12, color=MGRAY, align=PP_ALIGN.LEFT, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — IODB VALIDATION: HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "IODB Validation — How It Works",
    "Step-by-step validation logic from raw IODB to error report")

steps = [
    (BLUE,   "STEP 1: Load & Parse",
     "IODB uploaded — auto-detect sheet name, header row, and column structure — "
     "DataFrame built in memory with all rows and columns preserved"),
    (MID,    "STEP 2: Schema Check",
     "Verify mandatory columns are present (TAG NO, SERVICE, SIGNAL TYPE, JB, P&ID) — "
     "Flag any missing or misnamed columns as schema-level errors"),
    (rgb(0x1B,0x7A,0x4E), "STEP 3: Row-Level Rules",
     "For every data row: check for empty mandatory fields, duplicate tag numbers, "
     "invalid IO types, out-of-range values, and format violations"),
    (ORANGE, "STEP 4: Cross-Reference Checks",
     "Cross-validate field relationships: tag appears in loop but not cable schedule, "
     "JB assigned but no cable number, signal type conflicts with IO type"),
    (NAVY,   "STEP 5: Report Generation",
     "Compile all findings — write structured Excel report with: "
     "Sheet, Row, Column, Error Type, Severity (Critical / Warning / Info), Description"),
]
for i, (c, title, desc) in enumerate(steps):
    y = 1.28 + i*1.2
    rect(sl, 0.35, y, 12.65, 1.08, WHITE)
    rect(sl, 0.35, y, 0.55,  1.08, c)
    txt(sl, str(i+1), 0.37, y+0.26, 0.5, 0.5,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    txt(sl, title, 1.05, y+0.08, 4.5, 0.4,
        size=13, bold=True, color=NAVY, name="Calibri")
    txt(sl, desc,  1.05, y+0.52, 11.7, 0.45,
        size=11.5, color=MGRAY, name="Calibri")
    if i < 4:
        txt(sl, "v", 6.58, y+1.1, 0.5, 0.3,
            size=12, color=BLUE, align=PP_ALIGN.CENTER)

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — IODB VALIDATION: ERROR TYPES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "IODB Validation — Error Types & Report Format",
    "Every issue captured with location, type, severity, and plain-English description")

rect(sl, 0.35, 1.28, 12.65, 0.4, NAVY)
for hx, hw, ht in [(0.45,2.0,"Error Category"),(2.55,1.8,"Severity"),
                    (4.45,4.5,"Examples"),(9.05,3.85,"Business Risk")]:
    txt(sl, ht, hx, 1.31, hw, 0.34, size=12, bold=True,
        color=WHITE, name="Calibri")

rows_data = [
    (WHITE,   "Missing Mandatory Fields", "Critical",
     "TAG NO empty, SERVICE blank, SIGNAL TYPE not filled",
     "Document generation produces TBA everywhere; unacceptable for IFC"),
    (OFF_W,   "Duplicate Tag Numbers",    "Critical",
     "Same TAG NO appears on two rows in the IODB",
     "Two instruments with same tag — wiring conflict, commissioning error"),
    (WHITE,   "Invalid Signal/IO Type",   "Warning",
     "Signal type not in approved list (e.g. 'ANALOG OUT' spelled wrong)",
     "DCS I/O assignment fails; loop wiring sheet mismatch"),
    (OFF_W,   "Cross-Reference Mismatch","Warning",
     "Tag in Loop Input not found in Cable Schedule section of IODB",
     "Missing cable entries discovered late — delays field work"),
    (WHITE,   "Format Violations",        "Info",
     "JB number format incorrect, P&ID reference format wrong",
     "Document control rejections; non-conformance at client review"),
    (OFF_W,   "Orphaned JB References",   "Info",
     "JB assigned to tag but JB has no corresponding schedule entry",
     "Field installation cannot proceed without correct JB schedule"),
]
for ri, (bg, cat, sev, ex, risk) in enumerate(rows_data):
    y = 1.68 + ri*0.78
    rect(sl, 0.35, y, 12.65, 0.72, bg)
    txt(sl, cat,  0.45, y+0.06, 2.0, 0.58, size=11, bold=True, color=NAVY, name="Calibri")
    txt(sl, sev,  2.55, y+0.06, 1.8, 0.58, size=11, color=DGRAY, name="Calibri")
    txt(sl, ex,   4.45, y+0.06, 4.45, 0.58, size=10.5, color=MGRAY, name="Calibri")
    txt(sl, risk, 9.05, y+0.06, 3.75, 0.58, size=10.5, color=MGRAY, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — IODB VALIDATION: SAMPLE REPORT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "IODB Validation — Sample Report Output",
    "Structured Excel log — exact location, clear description, actionable guidance")

rect(sl, 0.35, 1.28, 12.65, 0.42, NAVY)
cols_r = [(0.45,0.9,"Row #"),(1.45,1.4,"Sheet"),
          (2.95,1.8,"Column"),(4.85,2.2,"Error Type"),
          (7.15,2.4,"Severity"),(9.65,3.2,"Description")]
for cx, cw, ch in cols_r:
    txt(sl, ch, cx, 1.31, cw, 0.36, size=12, bold=True, color=WHITE, name="Calibri")

sample_rows = [
    (WHITE,  "14",  "IODB","SERVICE",      "Missing Mandatory Field","Critical",
     "Cell is blank. SERVICE is required for datasheet generation."),
    (OFF_W,  "27",  "IODB","TAG NO",        "Duplicate Tag Number",  "Critical",
     "200-PT-20102 appears on rows 27 and 43. Tag numbers must be unique."),
    (WHITE,  "35",  "IODB","SIGNAL TYPE",   "Invalid Enumeration",   "Warning",
     "'ANLAOG IN' is not a recognised signal type. Did you mean 'ANALOG IN'?"),
    (OFF_W,  "52",  "IODB","JUNCTION BOX",  "Cross-Ref Mismatch",    "Warning",
     "Tag 200-FIT-20105 references JB-03 but no cable entry found for JB-03."),
    (WHITE,  "61",  "IODB","P&ID NO",       "Format Violation",      "Info",
     "P&ID reference '1234-E-5678-001' does not match expected format."),
    (OFF_W,  "78",  "IODB","LOOP NO",       "Empty Recommended Field","Info",
     "LOOP NO is blank. Loop wiring sheet will be generated without loop reference."),
    (WHITE,  "91",  "IODB","IO TYPE",       "Inconsistency",         "Warning",
     "IO TYPE = 'DO' but SIGNAL TYPE = 'ANALOG IN'. These are contradictory."),
]
for ri, (bg, row, sheet, col_, etype, sev, desc) in enumerate(sample_rows):
    y = 1.7 + ri*0.73
    rect(sl, 0.35, y, 12.65, 0.66, bg)
    for val, cx, cw in [
        (row,   0.45, 0.9), (sheet, 1.45, 1.4), (col_, 2.95, 1.8),
        (etype, 4.85, 2.2), (sev,   7.15, 2.4), (desc, 9.65, 3.2)
    ]:
        txt(sl, val, cx, y+0.06, cw, 0.55, size=10, color=DGRAY, name="Calibri")

rect(sl, 0.35, 6.85, 12.65, 0.35, LGRAY)
txt(sl,
    "Report saved as  IODB_Validation_Report.xlsx  — shareable with project team and QA lead for resolution",
    0.5, 6.87, 12.4, 0.3, size=11, color=MGRAY, align=PP_ALIGN.CENTER, name="Calibri")
foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — KEY BENEFITS & ROI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Key Benefits & Return on Investment",
    "Quantified impact across a typical 200-instrument project")

stats_b = [
    (">95%",   "Time Saved",       "Document prep from weeks\nto minutes",        BLUE),
    ("~170",   "Hrs / Project",    "Saved per 100-instrument\nproject",            MID),
    ("~0%",    "Error Rate",       "Transcription errors\neliminated",             GREEN),
    ("100%",   "Template Fidelity","All formats, formulas &\nshapes preserved",   ORANGE),
    ("Free",   "Re-Runs",          "Re-generate any time the\nIODB changes",       NAVY),
]
for i, (val, lbl, sub, c) in enumerate(stats_b):
    stat_card(sl, 0.32 + i*2.55, 1.28, 2.42, 1.52, val, lbl, sub, c)

rect(sl, 0.35, 2.95, 12.65, 0.42, NAVY)
for hx2, hw2, ht2 in [(0.45,3.2,"Deliverable"),
                       (3.75,2.8,"Manual Effort (100 instr)"),
                       (6.65,2.8,"With This Tool"),
                       (9.55,3.3,"Saving per Project")]:
    txt(sl, ht2, hx2, 2.98, hw2, 0.35, size=12, bold=True, color=WHITE, name="Calibri")

tbl_data = [
    (WHITE,  "Instrument List",  "3-4 hours",    "< 5 seconds",   "~3.5 hours"),
    (OFF_W,  "I/O List",         "4-6 hours",    "< 5 seconds",   "~5 hours"),
    (WHITE,  "Data Sheets",      "50-80 hours",  "< 1 hour",      "~65 hours"),
    (OFF_W,  "Cable Schedule",   "16-24 hours",  "< 0.5 hours",   "~20 hours"),
    (WHITE,  "Loop Wiring",      "50-100 hours", "< 1 hour",      "~75 hours"),
    (OFF_W,  "IODB Validation",  "Manual audit (if done at all)","< 10 sec","Defect prevention"),
    (LGRAY,  "TOTAL",            "~120-200 hours","~3 hours",     "~170 hours saved"),
]
for ri, (bg, doc, m, tool, sav) in enumerate(tbl_data):
    y = 3.37 + ri*0.54
    rect(sl, 0.35, y, 12.65, 0.5, bg)
    bold_row = (ri == 6)
    clr = NAVY if bold_row else DGRAY
    for val, cx, cw in [(doc,0.45,3.2),(m,3.75,2.8),(tool,6.65,2.8),(sav,9.55,3.3)]:
        txt(sl, val, cx, y+0.06, cw, 0.4, size=11,
            bold=bold_row, color=clr, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DEMO WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Live Demo — Walkthrough Steps",
    "How to demonstrate the tool in 8 minutes for a management audience")

demo_steps = [
    (BLUE,   "1. Open the App",
     "Navigate to the Streamlit URL in any browser — no login, no install.\n"
     "Show the clean header, sidebar navigation, and module list."),
    (MID,    "2. Upload IODB",
     "Upload IODB.xlsx once.\n"
     "Show the green 'Loaded Files' indicator — file persists across all modules."),
    (rgb(0x1B,0x7A,0x4E), "3. Instrument List",
     "Select 5-6 key columns, click Generate, download instantly.\n"
     "Open the file to show correct data and formatting."),
    (ORANGE, "4. Data Sheets",
     "Upload Datasheet template, select 3 tags, Generate, download ZIP.\n"
     "Open one .xlsx — show all fields filled, formulas intact."),
    (NAVY,   "5. Cable Schedule",
     "Upload CS template, click Generate, open output.\n"
     "Show natural-sorted JBs, SPARE rows, and preserved formatting."),
    (rgb(0x6B,0x35,0xAF), "6. IODB Validation",
     "Click Validate — completes in < 10 seconds.\n"
     "Open IODB_Validation_Report.xlsx — walk through a Critical error row."),
    (RED,    "7. Highlight a Fix",
     "Fix the flagged error in the IODB, re-upload, re-validate.\n"
     "Show the error disappears — close the feedback loop live."),
    (GOLD,   "8. Summary",
     "Recap: same IODB, 5 documents + validation report — all in under 5 minutes.\n"
     "Ask: 'How long would this take manually?'"),
]
for i, (c, title, desc) in enumerate(demo_steps):
    row = i // 4; col = i % 4
    x = 0.35 + col*3.25; y = 1.28 + row*2.85
    rect(sl, x, y, 3.1, 2.62, WHITE)
    rect(sl, x, y, 3.1, 0.44, c)
    txt(sl, title, x+0.12, y+0.06, 2.88, 0.36,
        size=12, bold=True, color=WHITE, name="Calibri")
    txt(sl, desc, x+0.12, y+0.56, 2.88, 1.88,
        size=11, color=MGRAY, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Future Roadmap & Scalability",
    "The platform is designed to grow — new modules, cloud deployment, and AI integration")

phases = [
    (GREEN,  "Phase 1 — Complete", [
        "Instrument List Generator",
        "I/O List Generator",
        "Data Sheet Generator (batch ZIP)",
        "Cable Schedule Generator",
        "Loop Wiring Generator",
        "IODB Validation Engine",
        "Session-state file caching",
        "Robust edge-case handling",
    ]),
    (BLUE,   "Phase 2 — Next Quarter", [
        "Enhanced validation rule library",
        "PDF export option for datasheets",
        "Multi-template support per module",
        "Automatic revision tracking",
        "Inline IODB edit and re-validate",
        "Email / MS Teams delivery",
        "User activity log for audit trail",
    ]),
    (NAVY,   "Phase 3 — Strategic Vision", [
        "Cloud deployment (Streamlit Cloud / Azure)",
        "Role-based access control (RBAC)",
        "Integration with project database (SQL)",
        "IODB change-diff between revisions",
        "AI-assisted column mapping & suggestions",
        "Approval workflow with digital stamp",
        "Multi-project / multi-discipline support",
    ]),
]
for i, (c, title, items) in enumerate(phases):
    x = 0.35 + i*4.35
    rect(sl, x, 1.28, 4.15, 5.65, WHITE)
    rect(sl, x, 1.28, 4.15, 0.48, c)
    txt(sl, title, x+0.14, 1.33, 3.92, 0.4,
        size=13, bold=True, color=WHITE, name="Calibri")
    body = sl.shapes.add_textbox(
        Inches(x+0.14), Inches(1.9), Inches(3.92), Inches(4.8))
    tf = body.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(8)
        r = p.add_run(); r.text = item
        r.font.size = Pt(12.5); r.font.color.rgb = DGRAY; r.font.name = "Calibri"

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — TECHNICAL ROBUSTNESS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Built for the Real World — Technical Robustness",
    "Designed to handle inconsistent files, naming variations, and edge cases gracefully")

cases = [
    (BLUE,   "Auto Sheet Detection",
     "Tries exact, case-insensitive, then substring match. "
     "'AI-INST ', 'AI Instr', 'ai instr' all resolve correctly."),
    (MID,    "Header Row Auto-Detection",
     "Probes rows 0-10 per sheet; scored by tag-column presence "
     "and non-empty column count. Works with metadata-heavy files."),
    (GREEN,  "Fuzzy Column Matching",
     "Template labels matched to IODB headers via normalised "
     "substring comparison — no manual mapping ever required."),
    (ORANGE, "Content_Types.xml Repair",
     "Synthesises a valid [Content_Types].xml from the ZIP manifest "
     "when a template was repaired by Excel and is corrupt."),
    (NAVY,   "Formula Preservation",
     "Every generator checks for '=' prefix before writing. "
     "Calculated fields remain live and correct in every output."),
    (RED,    "Drawing / Shape Injection",
     "ZIP-level XML copy injects complex drawings (connectors, "
     "P&ID symbols) that openpyxl cannot round-trip natively."),
    (rgb(0x6B,0x35,0xAF), "Column Name Tolerance",
     "'Tag_nummber', 'Tag No.', 'TAG NUMBER' — all auto-resolve "
     "to the same tag column. No user action needed."),
    (GOLD,   "File Format Guard",
     "Old .xls (OLE2) files detected by magic bytes and rejected "
     "with a clear 'Please save as .xlsx' message — not a crash."),
]
for i, (c, title, desc) in enumerate(cases):
    row = i // 4; col = i % 4
    x = 0.35 + col*3.25; y = 1.28 + row*2.85
    rect(sl, x, y, 3.1, 2.62, WHITE)
    rect(sl, x, y, 3.1, 0.07, c)
    txt(sl, title, x+0.12, y+0.14, 2.88, 0.42,
        size=12, bold=True, color=NAVY, name="Calibri")
    txt(sl, desc, x+0.12, y+0.62, 2.88, 1.85,
        size=11, color=MGRAY, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CONCLUSION & BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_base(sl)
hdr(sl, "Conclusion — Business Impact Statement",
    "Instrumentation EZ delivers measurable, lasting value across every project phase")

rect(sl, 0.35, 1.28, 12.65, 0.72, NAVY)
txt(sl,
    "This platform fundamentally changes how instrumentation engineering documents are produced — "
    "faster, error-free, and production-ready from day one.",
    0.5, 1.33, 12.35, 0.6, size=14, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")

impacts = [
    (BLUE,  "Operational Impact",
     ["Document turnaround from days to minutes",
      "~170 engineering hours saved per project",
      "Re-generation on IODB change costs zero extra effort",
      "All 5 document types from a single data upload",
      "Zero manual formatting — template fidelity guaranteed"]),
    (GREEN, "Quality Impact",
     ["Transcription errors eliminated by design",
      "IODB Validation catches data issues at source",
      "TBA safety net — no blank cells in any output",
      "Consistent document style across all engineers",
      "Audit trail: validation report documents QA evidence"]),
    (ORANGE,"Strategic Impact",
     ["Frees senior engineers for value-added work",
      "Reduced rework cost at client review stage",
      "Platform scales to any project size instantly",
      "Foundation for future AI and cloud integration",
      "Competitive advantage in tender & execution speed"]),
]
for i, (c, title, bullets) in enumerate(impacts):
    x = 0.35 + i*4.35
    rect(sl, x, 2.1, 4.1, 4.75, WHITE)
    rect(sl, x, 2.1, 4.1, 0.48, c)
    txt(sl, title, x+0.14, 2.15, 3.9, 0.38,
        size=13, bold=True, color=WHITE, name="Calibri")
    body = sl.shapes.add_textbox(
        Inches(x+0.14), Inches(2.7), Inches(3.9), Inches(3.9))
    tf = body.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(8)
        r = p.add_run(); r.text = "- " + b
        r.font.size = Pt(12.5); r.font.color.rgb = DGRAY; r.font.name = "Calibri"

rect(sl, 0.35, 6.97, 12.65, 0.35, BLUE)
txt(sl,
    "Recommended Next Step: Approve platform for live project deployment "
    "and initiate Phase 2 enhancements",
    0.5, 6.99, 12.35, 0.3, size=12, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")

foot(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 0, 13.33, 0.08, BLUE)
rect(sl, 0, 7.42, 13.33, 0.08, BLUE)

acc2 = sl.shapes.add_shape(1, Inches(8.5), Inches(0), Inches(4.83), Inches(7.5))
acc2.fill.solid(); acc2.fill.fore_color.rgb = rgb(0x12,0x38,0x62)
acc2.line.fill.background()

txt(sl, "Thank You", 0.8, 0.9, 9.0, 1.1,
    size=52, bold=True, color=WHITE,
    align=PP_ALIGN.LEFT, name="Calibri Light")
rect(sl, 0.8, 2.1, 6.0, 0.06, BLUE)
txt(sl, "Instrumentation EZ — Key Takeaways",
    0.8, 2.28, 8.0, 0.5, size=18,
    color=rgb(0xA8,0xC8,0xEE), align=PP_ALIGN.LEFT, name="Calibri Light")

takeaways = [
    "6 integrated modules — document generation and data validation in one platform",
    "~170 engineering hours saved per 100-instrument project",
    "IODB Validation catches data errors before they reach any deliverable",
    "Zero transcription errors — all outputs flow directly from IODB",
    "Built for extensibility — cloud, AI, and workflow integration ready",
]
body = sl.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(7.8), Inches(2.8))
tf = body.text_frame; tf.word_wrap = True
first = True
for t in takeaways:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False; p.space_before = Pt(10)
    r = p.add_run(); r.text = "◆  " + t
    r.font.size = Pt(13.5); r.font.color.rgb = rgb(0xC5,0xD8,0xF0); r.font.name = "Calibri"

txt(sl,
    "Built with  Python · Streamlit · openpyxl  |  GitHub: AdithyaChoudhry/IEZ",
    0.8, 6.0, 8.0, 0.4, size=12,
    color=rgb(0x70,0x90,0xB0), align=PP_ALIGN.LEFT, name="Calibri")
txt(sl, "Questions & Discussion",
    0.8, 6.55, 8.0, 0.6, size=22, bold=True,
    color=WHITE, align=PP_ALIGN.LEFT, name="Calibri Light")


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "Instrumentation_EZ_Corporate_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
